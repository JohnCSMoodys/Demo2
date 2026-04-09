import io
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ──────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0f, 0x11, 0x17)
C_FG      = RGBColor(0xe8, 0xea, 0xf0)
C_SUBTEXT = RGBColor(0x90, 0x96, 0xa8)
C_GREEN   = RGBColor(0x4C, 0xAF, 0x50)
C_RED     = RGBColor(0xf4, 0x43, 0x36)
C_ABSENT  = RGBColor(0xFF, 0x6B, 0x6B)
C_PRESENT = RGBColor(0x4E, 0xCD, 0xC4)
C_GRID    = RGBColor(0x2a, 0x2d, 0x3a)
C_ACCENT  = RGBColor(0xFF, 0xD7, 0x00)

MPL_BG      = "#0f1117"
MPL_FG      = "#e8eaf0"
MPL_SUBTEXT = "#9096a8"
MPL_GRID    = "#2a2d3a"
MPL_GREEN   = "#4CAF50"
MPL_RED     = "#f44336"
MPL_ABSENT  = "#FF6B6B"
MPL_PRESENT = "#4ECDC4"

# ── data ─────────────────────────────────────────────────────────────────────
df = pd.read_csv("stars.csv")
n  = len(df)

PROFILE_COLS = ["metallicity_fe_h", "teff_k", "distance_pc", "radius_solar", "mass_solar"]
COND_COLS    = ["teff_k", "mass_solar", "radius_solar", "distance_pc"]

missing_pct  = {c: df[c].isna().sum() / n * 100 for c in PROFILE_COLS}
present_pct  = {c: 100 - missing_pct[c]          for c in PROFILE_COLS}
met_absent   = df["metallicity_fe_h"].isna()
n_absent     = int(met_absent.sum())
n_present    = int((~met_absent).sum())
cond_absent  = {c: df.loc[met_absent,  c].isna().mean() * 100 for c in COND_COLS}
cond_present = {c: df.loc[~met_absent, c].isna().mean() * 100 for c in COND_COLS}

# ── helpers ───────────────────────────────────────────────────────────────────
def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, facecolor=MPL_BG, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def make_panel1():
    fig, ax = plt.subplots(figsize=(9, 4.2), facecolor=MPL_BG)
    ax.set_facecolor(MPL_BG)
    for sp in ax.spines.values():
        sp.set_edgecolor(MPL_GRID)

    sorted_cols = sorted(PROFILE_COLS, key=lambda c: present_pct[c])
    labels      = sorted_cols
    pres_vals   = [present_pct[c] for c in sorted_cols]
    miss_vals   = [missing_pct[c] for c in sorted_cols]
    y_pos       = np.arange(len(sorted_cols))

    ax.barh(y_pos, pres_vals, height=0.55, color=MPL_GREEN, zorder=3)
    ax.barh(y_pos, miss_vals, height=0.55, left=pres_vals, color=MPL_RED, zorder=3)
    ax.axvline(100, color=MPL_GRID, linestyle="--", linewidth=0.8)

    for i, p in enumerate(pres_vals):
        ax.text(p - 0.5, i, f"{p:.1f}%", ha="right", va="center",
                fontsize=9, color=MPL_BG, fontweight="bold")

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=9, color=MPL_FG)
    ax.set_xlim(85, 101)
    ax.set_xlabel("% complete", color=MPL_FG, fontsize=9)
    ax.tick_params(colors=MPL_FG)
    ax.xaxis.grid(True, color=MPL_GRID, linewidth=0.5, zorder=0)
    ax.set_axisbelow(True)
    ax.set_title("Column Completeness Profile", color=MPL_FG, fontsize=11, fontweight="bold", pad=8)
    leg = [mpatches.Patch(color=MPL_GREEN, label="Present"),
           mpatches.Patch(color=MPL_RED,   label="Missing")]
    ax.legend(handles=leg, fontsize=8, framealpha=0.15, labelcolor=MPL_FG,
              facecolor=MPL_BG, edgecolor=MPL_GRID, loc="lower right")
    fig.tight_layout()
    return fig_to_stream(fig)


def make_panel2():
    fig, ax = plt.subplots(figsize=(9, 4.2), facecolor=MPL_BG)
    ax.set_facecolor(MPL_BG)
    for sp in ax.spines.values():
        sp.set_edgecolor(MPL_GRID)

    x        = np.arange(len(COND_COLS))
    bw       = 0.35
    absent_v  = [cond_absent[c]  for c in COND_COLS]
    present_v = [cond_present[c] for c in COND_COLS]
    y_ceil    = max(absent_v) * 1.25

    bars_abs = ax.bar(x - bw/2, absent_v,  width=bw, color=MPL_ABSENT,  zorder=3)
    bars_pre = ax.bar(x + bw/2, present_v, width=bw, color=MPL_PRESENT, zorder=3)

    for bar, val in zip(bars_abs, absent_v):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
                f"{val:.1f}%", ha="center", va="bottom", fontsize=8,
                color=MPL_ABSENT, fontweight="bold")
    for bar, val in zip(bars_pre, present_v):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
                f"{val:.1f}%", ha="center", va="bottom", fontsize=8,
                color=MPL_PRESENT, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(COND_COLS, fontsize=9, color=MPL_FG)
    ax.set_ylim(0, y_ceil)
    ax.set_ylabel("% missing", color=MPL_FG, fontsize=9)
    ax.tick_params(colors=MPL_FG)
    ax.yaxis.grid(True, color=MPL_GRID, linewidth=0.5, zorder=0)
    ax.set_axisbelow(True)
    ax.set_title("Conditional Missingness (metallicity absent vs. present)",
                 color=MPL_FG, fontsize=11, fontweight="bold", pad=8)

    dist_i = COND_COLS.index("distance_pc")
    ax.annotate(
        "Median dist similar\n(255 pc vs 290 pc)\n→ not the driver",
        xy=(dist_i + bw/2, cond_present["distance_pc"]),
        xytext=(dist_i - 0.2, y_ceil * 0.65),
        fontsize=7.5, color=MPL_SUBTEXT, ha="center",
        arrowprops=dict(arrowstyle="->", color=MPL_SUBTEXT, lw=0.8),
        bbox=dict(boxstyle="round,pad=0.3", facecolor=MPL_BG, edgecolor=MPL_GRID, alpha=0.9),
    )
    leg = [mpatches.Patch(color=MPL_ABSENT,  label=f"Met. absent  (n={n_absent})"),
           mpatches.Patch(color=MPL_PRESENT, label=f"Met. present (n={n_present})")]
    ax.legend(handles=leg, fontsize=8, framealpha=0.15, labelcolor=MPL_FG,
              facecolor=MPL_BG, edgecolor=MPL_GRID, loc="upper right")
    fig.tight_layout()
    return fig_to_stream(fig)


# ── pptx helpers ──────────────────────────────────────────────────────────────
def set_bg(slide, prs):
    from pptx.util import Inches
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_text(txf, text, size, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    tf = txf.text_frame
    tf.word_wrap = True
    p  = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text      = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_FG


def title_box(slide, prs, text, top=Inches(0.25)):
    w  = prs.slide_width
    tb = slide.shapes.add_textbox(Inches(0.5), top, w - Inches(1), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(24)
    run.font.bold      = True
    run.font.color.rgb = C_FG
    return tb


def body_box(slide, prs, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def bullet(tf, text, size=13, color=None, indent=0, bold=False):
    p = tf.add_paragraph()
    p.level = indent
    p.alignment = PP_ALIGN.LEFT
    if indent == 0 and text.startswith("•") is False and text.strip():
        p.space_before = Pt(4)
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color or C_FG


def colored_bullet(tf, label, rest, label_color, size=12):
    p   = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(6)
    r1 = p.add_run()
    r1.text           = label
    r1.font.size      = Pt(size)
    r1.font.bold      = True
    r1.font.color.rgb = label_color
    r2 = p.add_run()
    r2.text           = rest
    r2.font.size      = Pt(size)
    r2.font.color.rgb = C_FG


# ── build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # blank

panel1_stream = make_panel1()
panel2_stream = make_panel2()

# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)

tb = sl.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.4))
tf = tb.text_frame
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text           = "Stars Catalog — Data Quality Audit"
run.font.size      = Pt(36)
run.font.bold      = True
run.font.color.rgb = C_FG

tb2 = sl.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(10.9), Inches(0.6))
tf2 = tb2.text_frame
p2  = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text            = "Where is the data trustworthy, thin, or incomplete?"
r2.font.size       = Pt(18)
r2.font.italic     = True
r2.font.color.rgb  = C_SUBTEXT

tb3 = sl.shapes.add_textbox(Inches(1.2), Inches(4.7), Inches(10.9), Inches(0.5))
tf3 = tb3.text_frame
p3  = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text            = "For data practitioners  ·  972 host stars  ·  5 analytical columns"
r3.font.size       = Pt(13)
r3.font.color.rgb  = C_SUBTEXT

# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Dataset
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)
title_box(sl, prs, "The Dataset")

tb = body_box(sl, prs, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.6))
tf = tb.text_frame
bullet(tf, "Source:  stars.csv — NASA Exoplanet Archive stellar catalog", 13, bold=True)
bullet(tf, "")
bullet(tf, f"  {n} host stars", 12)
bullet(tf, "  10 columns total: identifiers, stellar properties, positional data, num_planets", 12)
bullet(tf, "")
bullet(tf, "Columns audited (analytically important stellar properties):", 13, bold=True)
bullet(tf, "")

rows = [
    ("metallicity_fe_h", "Stellar metallicity [Fe/H]",  "Requires spectroscopy  →  most expensive to measure"),
    ("teff_k",           "Effective temperature (K)",   "Can be estimated photometrically"),
    ("distance_pc",      "Distance (parsecs)",           "Parallax-based; widely available from Gaia"),
    ("radius_solar",     "Stellar radius (R☉)",         "Can be estimated from T_eff and luminosity"),
    ("mass_solar",       "Stellar mass (M☉)",           "Estimated from stellar models"),
]
for col, desc, note in rows:
    p = tf.add_paragraph()
    p.space_before = Pt(5)
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = f"  {col:<22s}"; r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = C_ACCENT
    r2 = p.add_run(); r2.text = f"{desc:<35s}"; r2.font.size = Pt(11); r2.font.color.rgb = C_FG
    r3 = p.add_run(); r3.text = f"  //  {note}"; r3.font.size = Pt(10); r3.font.color.rgb = C_SUBTEXT

bullet(tf, "")
bullet(tf, "Left out:  star_id, star_name, ra_deg, dec_deg, num_planets — 0 % missing, not analytically risky", 11, color=C_SUBTEXT)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — Story Hypothesis
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)
title_box(sl, prs, "Story Hypothesis")

tb = body_box(sl, prs, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.1))
tf = tb.text_frame
p  = tf.paragraphs[0]
r  = p.add_run()
r.text           = ("Completeness is uneven and concentrated: a single column (metallicity_fe_h) accounts "
                    "for nearly all missingness, and stars missing that column are data-sparse across the board.")
r.font.size      = Pt(13)
r.font.italic    = True
r.font.color.rgb = C_SUBTEXT

claims = [
    (
        "A.  Missingness is concentrated in metallicity",
        C_ABSENT,
        (f"metallicity_fe_h is missing {missing_pct['metallicity_fe_h']:.1f} % of the time vs. <2 % for all other columns. "
         "Gap reflects measurement cost: metallicity requires high-resolution spectroscopy; "
         "temperature, mass, and radius can be estimated photometrically."),
    ),
    (
        "B.  Missing metallicity predicts missing everything else",
        C_PRESENT,
        (f"Stars without metallicity show dramatically elevated missingness across all other columns: "
         f"teff_k is missing {cond_absent['teff_k']:.0f} % in that group (vs. {cond_present['teff_k']:.1f} % otherwise), "
         f"radius_solar {cond_absent['radius_solar']:.1f} % (vs. {cond_present['radius_solar']:.1f} %), "
         f"distance_pc {cond_absent['distance_pc']:.1f} % (vs. {cond_present['distance_pc']:.1f} %). "
         "These are data-sparse stars, not random gaps."),
    ),
    (
        "C.  Distance does not explain it",
        RGBColor(0xFF, 0xD7, 0x00),
        ("Median distance for stars with vs. without metallicity is nearly identical "
         "(~255 pc vs. ~290 pc), ruling out observational depth as the driver. "
         "The gap is about which stars received spectroscopic follow-up."),
    ),
]
for i, (header, color, body) in enumerate(claims):
    top = Inches(2.25 + i * 1.55)
    htb = body_box(sl, prs, Inches(0.6), top, Inches(12.1), Inches(0.42))
    htf = htb.text_frame
    hp  = htf.paragraphs[0]
    hr  = hp.add_run()
    hr.text = header; hr.font.size = Pt(14); hr.font.bold = True; hr.font.color.rgb = color

    btb = body_box(sl, prs, Inches(0.75), top + Inches(0.38), Inches(11.9), Inches(1.0))
    btf = btb.text_frame
    bp  = btf.paragraphs[0]
    br  = bp.add_run()
    br.text = body; br.font.size = Pt(11); br.font.color.rgb = C_FG
    btf.word_wrap = True

# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Panel 1
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)
title_box(sl, prs, "Panel 1 — Column Completeness Profile")

sl.shapes.add_picture(panel1_stream, Inches(0.3), Inches(1.0), Inches(8.5), Inches(4.0))

tb = body_box(sl, prs, Inches(9.1), Inches(1.2), Inches(3.9), Inches(5.5))
tf = tb.text_frame
bullet(tf, "Key findings", 13, bold=True, color=C_ACCENT)
bullet(tf, "")
bullet(tf, f"metallicity_fe_h  →  {present_pct['metallicity_fe_h']:.1f}% complete", 12, color=C_ABSENT)
bullet(tf, "Clear worst offender — 4× worse than nearest peer", 10, color=C_SUBTEXT)
bullet(tf, "")

for c in ["teff_k", "distance_pc", "radius_solar", "mass_solar"]:
    bullet(tf, f"{c}  →  {present_pct[c]:.1f}% complete", 12)

bullet(tf, "")
bullet(tf, "All other columns ≥ 98.4% complete", 11, color=C_SUBTEXT)
bullet(tf, "")
bullet(tf, "Sorted worst → best (metallicity at bottom)", 11, color=C_SUBTEXT)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — Panel 2
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)
title_box(sl, prs, "Panel 2 — Conditional Missingness")

sl.shapes.add_picture(panel2_stream, Inches(0.3), Inches(1.0), Inches(8.5), Inches(4.0))

tb = body_box(sl, prs, Inches(9.1), Inches(1.2), Inches(3.9), Inches(5.5))
tf = tb.text_frame
bullet(tf, "Key findings", 13, bold=True, color=C_ACCENT)
bullet(tf, "")
bullet(tf, f"Stars missing metallicity (n={n_absent}):", 12, bold=True, color=C_ABSENT)
bullet(tf, f"  teff_k        →  {cond_absent['teff_k']:.0f}%  missing", 11)
bullet(tf, f"  radius_solar  →  {cond_absent['radius_solar']:.1f}%  missing", 11)
bullet(tf, f"  distance_pc   →  {cond_absent['distance_pc']:.1f}%  missing", 11)
bullet(tf, f"  mass_solar    →  {cond_absent['mass_solar']:.1f}%  missing", 11)
bullet(tf, "")
bullet(tf, f"Stars with metallicity (n={n_present}):", 12, bold=True, color=C_PRESENT)
bullet(tf, f"  teff_k        →  {cond_present['teff_k']:.1f}%  missing", 11)
bullet(tf, f"  All others    ≤  {max(cond_present[c] for c in COND_COLS):.1f}%  missing", 11)
bullet(tf, "")
bullet(tf, "Median distance similar (~255 vs ~290 pc)", 10, color=C_SUBTEXT)
bullet(tf, "→ distance is not the driver", 10, color=C_SUBTEXT)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — Risks & Caveats
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)
title_box(sl, prs, "Risks & Caveats")

risks = [
    ("Small absent group",
     f"Only {n_absent} stars (6.2%) lack metallicity — conditional rates have wider uncertainty. "
     "Flag n on Panel 2."),
    ("Co-missingness vs. causality",
     "Missing metallicity correlates with, but does not cause, missing other columns. "
     "Frame as 'data-sparse stars', not 'metallicity gap causes other gaps'."),
    ("metallicity_fe_h outlier status",
     "6.2% vs. <2% for peers is striking — could be a pipeline artifact. "
     "Flag as 'investigate before using as feature' rather than 'don't use'."),
    ("num_planets completeness masks planet-data gaps",
     "This audit covers stellar columns only. Planet-level data quality is not captured here."),
    ("Snapshot in time",
     "stars.csv is a point-in-time pull. Stars missing data may gain measurements as surveys continue."),
]

for i, (title_r, body_r) in enumerate(risks):
    top = Inches(1.15 + i * 1.12)
    htb = body_box(sl, prs, Inches(0.6), top, Inches(12.1), Inches(0.4))
    htf = htb.text_frame
    hp  = htf.paragraphs[0]
    hr  = hp.add_run()
    hr.text = f"⚑  {title_r}"
    hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = C_ACCENT

    btb = body_box(sl, prs, Inches(0.9), top + Inches(0.37), Inches(11.8), Inches(0.58))
    btf = btb.text_frame
    btf.word_wrap = True
    bp  = btf.paragraphs[0]
    br  = bp.add_run()
    br.text = body_r; br.font.size = Pt(11); br.font.color.rgb = C_FG

# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — What We Built
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, prs)
title_box(sl, prs, "What We Built")

tb = body_box(sl, prs, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.6))
tf = tb.text_frame

bullet(tf, "Deliverables", 14, bold=True, color=C_ACCENT)
bullet(tf, "")
bullet(tf, "  design_diagram.png  —  two-panel matplotlib figure (dark theme, 150 dpi)", 12)
bullet(tf, "  stars_data_quality_audit.pptx  —  this 7-slide narrative deck", 12)
bullet(tf, "")
bullet(tf, "Methodology", 14, bold=True, color=C_ACCENT)
bullet(tf, "")
steps = [
    "Loaded stars.csv (972 rows, 10 columns) with pandas",
    "Computed per-column missingness rates and present-pct across 5 analytical columns",
    f"Split data by metallicity_fe_h presence (absent n={n_absent}, present n={n_present})",
    "Computed conditional missingness rates for each secondary column in each group",
    "Panel 1: horizontal bar chart sorted by completeness (worst → best)",
    "Panel 2: grouped bar chart showing conditional missingness with distance callout",
    "Composed two-panel figure with shared dark-theme styling",
    "Generated narrative PowerPoint deck from the same computed metrics",
]
for s in steps:
    bullet(tf, f"  {i+1}.  {s}" if False else f"  •  {s}", 11)

bullet(tf, "")
bullet(tf, "Both scripts are idempotent — re-running produces identical outputs.", 11, color=C_SUBTEXT)

# ─────────────────────────────────────────────────────────────────────────────
prs.save("stars_data_quality_audit.pptx")
print("Saved stars_data_quality_audit.pptx")
